#!/usr/bin/env python3
"""extract.py の検証。依存なしで動く（python tests/test_extract.py）。

実際に docx / pptx / xlsx / png を生成して往復させ、テキストと画像が取り出せることを確認する。
失敗すると終了コード 1。
"""
from __future__ import annotations

import io
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import extract  # noqa: E402

FAILED: list[str] = []


def check(cond: bool, label: str) -> None:
    print(("  ok   " if cond else "  FAIL ") + label)
    if not cond:
        FAILED.append(label)


def png_bytes(w: int = 40, h: int = 30, color=(200, 30, 30)) -> bytes:
    from PIL import Image as PILImage

    buf = io.BytesIO()
    PILImage.new("RGB", (w, h), color).save(buf, format="PNG")
    return buf.getvalue()


# ---------------- sniff ----------------
def test_sniff() -> None:
    print("sniff")
    check(extract.sniff("a.docx") == "docx", "docx by extension")
    check(extract.sniff("a.pptx") == "pptx", "pptx by extension")
    check(extract.sniff("a.xlsx") == "xlsx", "xlsx by extension")
    check(extract.sniff("a.PDF") == "pdf", "pdf is case-insensitive")
    check(extract.sniff("a.md") == "text", "markdown is text")
    check(extract.sniff("a.PNG") == "image", "png is image")
    check(extract.sniff("noext", "image/jpeg") == "image", "falls back to mimetype")
    check(extract.sniff("weird.bin") == "unknown", "unknown stays unknown")


# ---------------- docx ----------------
def test_docx() -> None:
    print("docx")
    import docx

    d = docx.Document()
    d.add_paragraph("ウェブサイト改善提案")
    d.add_paragraph("H｜課題と改善仮説")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "KPI"
    t.cell(0, 1).text = "目標値"
    t.cell(1, 0).text = "申込フォーム遷移数"
    t.cell(1, 1).text = "月150件"
    buf = io.BytesIO()
    d.save(buf)
    data = buf.getvalue()

    r = extract.extract_text(data, "report.docx")
    check(r["kind"] == "docx", "kind is docx")
    check("ウェブサイト改善提案" in r["text"], "paragraph text extracted")
    check("H｜課題と改善仮説" in r["text"], "japanese pipe char survives")
    check("申込フォーム遷移数 | 月150件" in r["text"], "table row extracted")
    check(r["note"] == "", "no note on a clean docx")
    check(r["truncated"] is False, "not truncated")

    r2 = extract.extract_text(data, "report.docx", max_chars=10)
    check(r2["truncated"] is True and r2["chars"] == 10, "max_chars truncates")
    check("打ち切り" in r2["note"], "truncation is reported in note")


def test_docx_images() -> None:
    print("docx images")
    import docx

    d = docx.Document()
    d.add_paragraph("図あり")
    d.add_picture(io.BytesIO(png_bytes()))
    buf = io.BytesIO()
    d.save(buf)

    imgs, note = extract.extract_images(buf.getvalue(), "report.docx")
    check(len(imgs) == 1, f"one image from docx (got {len(imgs)})")
    check(imgs and imgs[0]["format"] in ("png", "jpeg"), "image has a usable format")
    check(imgs and len(imgs[0]["data"]) > 0, "image bytes are non-empty")


# ---------------- pptx ----------------
def test_pptx() -> None:
    print("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = "課題3｜ウェブサイト改善提案"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "課題4｜ウェブマーケティング提案"
    s2.shapes.add_picture(io.BytesIO(png_bytes(60, 40)), Inches(1), Inches(2))
    s2.notes_slide.notes_text_frame.text = "発表メモ"
    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()

    r = extract.extract_text(data, "teian.pptx")
    check(r["kind"] == "pptx", "kind is pptx")
    check("--- スライド 1 ---" in r["text"], "slide 1 marker present")
    check("--- スライド 2 ---" in r["text"], "slide 2 marker present")
    check("課題3｜ウェブサイト改善提案" in r["text"], "slide 1 title extracted")
    check("[ノート] 発表メモ" in r["text"], "speaker notes extracted")

    imgs, _ = extract.extract_images(data, "teian.pptx")
    check(len(imgs) == 1, f"one image from pptx (got {len(imgs)})")
    check(imgs and imgs[0]["name"].startswith("slide2_"),
          f"image is tagged with its slide number (got {imgs[0]['name'] if imgs else None})")


# ---------------- xlsx ----------------
def test_xlsx() -> None:
    print("xlsx")
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "KPI設定シート"
    ws.append(["ステージ", "施策", "KPI"])
    ws.append(["認知", "検索広告", "表示回数"])
    ws2 = wb.create_sheet("メモ")
    ws2.append(["備考", "テスト"])
    buf = io.BytesIO()
    wb.save(buf)

    r = extract.extract_text(buf.getvalue(), "sheet.xlsx")
    check(r["kind"] == "xlsx", "kind is xlsx")
    check("--- シート: KPI設定シート ---" in r["text"], "sheet name emitted")
    check("--- シート: メモ ---" in r["text"], "second sheet emitted")
    check("認知\t検索広告\t表示回数" in r["text"], "row values are tab-joined")


# ---------------- text / image / errors ----------------
def test_text_and_image() -> None:
    print("text / image / errors")
    r = extract.extract_text("見出し\n本文".encode("utf-8"), "note.md")
    check(r["kind"] == "text" and "本文" in r["text"], "utf-8 text decoded")

    r = extract.extract_text("日本語".encode("cp932"), "sjis.txt")
    check("日本語" in r["text"], "cp932 fallback decoding works")

    r = extract.extract_text(png_bytes(), "shot.png")
    check(r["kind"] == "image" and r["text"] == "", "image yields no text")
    check("read_submission_images" in r["note"], "image note points at the image tool")

    r = extract.extract_text(b"\x00\x01\x02", "mystery.bin")
    check(r["kind"] == "unknown" and "未対応" in r["note"], "unknown format is reported")

    r = extract.extract_text(b"not a real docx", "broken.docx")
    check(r["text"] == "" and "失敗" in r["note"], "corrupt docx reports, does not raise")

    imgs, note = extract.extract_images(b"nope", "broken.pptx")
    check(imgs == [] and note != "", "corrupt pptx reports, does not raise")

    imgs, note = extract.extract_images(b"x", "notes.txt")
    check(imgs == [] and "取り出せません" in note, "text file has no images")


def test_image_downscale() -> None:
    print("image downscale")
    from PIL import Image as PILImage

    big = png_bytes(3000, 2000, (10, 120, 200))
    imgs, _ = extract.extract_images(big, "big.png", max_px=1568)
    check(len(imgs) == 1, "large png accepted")
    im = PILImage.open(io.BytesIO(imgs[0]["data"]))
    check(max(im.size) <= 1568, f"downscaled to <=1568px (got {im.size})")
    check(len(imgs[0]["data"]) <= 3_500_000, "fits the per-image byte budget")

    small = png_bytes(40, 30)
    imgs, _ = extract.extract_images(small, "small.png", max_px=1568)
    im = PILImage.open(io.BytesIO(imgs[0]["data"]))
    check(im.size == (40, 30), "small image is not upscaled")


def test_image_limit() -> None:
    print("image limit")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for i in range(4):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_picture(io.BytesIO(png_bytes(30 + i, 20, (i * 40, 60, 90))),
                             Inches(1), Inches(1))
    buf = io.BytesIO()
    prs.save(buf)

    imgs, note = extract.extract_images(buf.getvalue(), "many.pptx", limit=2)
    check(len(imgs) == 2, f"limit is honoured (got {len(imgs)})")
    check("4枚" in note and "2枚" in note, f"note states how many were dropped ({note})")


if __name__ == "__main__":
    for fn in (test_sniff, test_docx, test_docx_images, test_pptx, test_xlsx,
               test_text_and_image, test_image_downscale, test_image_limit):
        fn()
    print()
    if FAILED:
        print(f"FAILED {len(FAILED)}:")
        for f in FAILED:
            print("  - " + f)
        sys.exit(1)
    print("all checks passed")
