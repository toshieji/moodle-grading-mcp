#!/usr/bin/env python3
"""Text and image extraction from submitted files / 提出ファイルの本文・画像抽出.

Pure functions with no Moodle dependency, so they can be unit-tested on their own.
Moodle 非依存の純粋関数。単体テストできるよう server.py から分離している。

Supported / 対応形式:
  .docx  python-docx   paragraphs + tables
  .pptx  python-pptx   per-slide shape text + tables + speaker notes
  .xlsx  openpyxl      per-sheet cell values (cached values, not formulas)
  .pdf   pypdf         per-page text
  text   -             .txt .md .csv .tsv .json .html
  image  Pillow        .png .jpg .gif .webp .bmp (returned as image blocks, not text)

Images embedded in docx/pptx/xlsx/pdf are extracted too, so a grader model can read
figures, screenshots and wireframes that carry no extractable text.
docx/pptx/xlsx/pdf に埋め込まれた画像も取り出せる。図・スクリーンショット・ワイヤー
フレームなど、テキストとして取り出せない情報をモデルに見せて書き起こさせるため。
"""
from __future__ import annotations

import io
import os
import zipfile

# 抽出テキストの既定上限（文字）。呼び出し側で上書き可。
DEFAULT_MAX_CHARS = 800_000
# 画像1枚の長辺の既定上限(px)。これ以上は縮小する。
DEFAULT_MAX_PX = 1568
# 画像1枚あたりの目標バイト数。超えたら JPEG へ落として再エンコードする。
_IMAGE_TARGET_BYTES = 3_500_000

_TEXT_EXT = {".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".html", ".htm", ".xml", ".log"}
_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"}


def sniff(filename: str, mimetype: str = "") -> str:
    """Return a coarse kind: docx / pptx / xlsx / pdf / text / image / unknown."""
    ext = os.path.splitext(filename or "")[1].lower()
    if ext == ".docx":
        return "docx"
    if ext == ".pptx":
        return "pptx"
    if ext in (".xlsx", ".xlsm"):
        return "xlsx"
    if ext == ".pdf":
        return "pdf"
    if ext in _TEXT_EXT:
        return "text"
    if ext in _IMAGE_EXT:
        return "image"
    mt = (mimetype or "").lower()
    if mt.startswith("image/"):
        return "image"
    if mt.startswith("text/") or "json" in mt:
        return "text"
    if "pdf" in mt:
        return "pdf"
    if "wordprocessingml" in mt:
        return "docx"
    if "presentationml" in mt:
        return "pptx"
    if "spreadsheetml" in mt:
        return "xlsx"
    return "unknown"


def _decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp932", "euc-jp"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _docx_text(data: bytes) -> str:
    import docx  # python-docx

    d = docx.Document(io.BytesIO(data))
    out: list[str] = []
    for p in d.paragraphs:
        t = (p.text or "").strip()
        if t:
            out.append(t)
    for i, tbl in enumerate(d.tables, 1):
        out.append(f"--- 表 {i} ---")
        for row in tbl.rows:
            cells = [(c.text or "").strip().replace("\n", " ") for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
    return "\n".join(out)


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    for n, slide in enumerate(prs.slides, 1):
        out.append(f"--- スライド {n} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = (shape.text_frame.text or "").strip()
                if t:
                    out.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [(c.text or "").strip().replace("\n", " ") for c in row.cells]
                    if any(cells):
                        out.append(" | ".join(cells))
        try:
            if slide.has_notes_slide:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
                if note:
                    out.append(f"[ノート] {note}")
        except Exception:
            pass
    return "\n".join(out)


def _xlsx_text(data: bytes) -> str:
    import openpyxl

    # data_only=True はキャッシュ済みの計算結果を返す。Excel で一度も開かれていない
    # ファイルでは数式セルが None になり得る（その旨を末尾に注記する）。
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out: list[str] = []
    empty_formula_cells = False
    try:
        for ws in wb.worksheets:
            out.append(f"--- シート: {ws.title} ---")
            for row in ws.iter_rows(values_only=True):
                vals = ["" if v is None else str(v).strip() for v in row]
                if any(vals):
                    out.append("\t".join(vals).rstrip())
                elif row and all(v is None for v in row):
                    empty_formula_cells = True
    finally:
        wb.close()
    if empty_formula_cells:
        out.append(
            "[注記] 空欄が多い場合、数式の計算結果が保存されていない可能性があります"
            "（Excel で開いて保存し直すと値が入ります）。"
        )
    return "\n".join(out)


def _pdf_text(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    out: list[str] = []
    for n, page in enumerate(reader.pages, 1):
        out.append(f"--- ページ {n} ---")
        try:
            out.append((page.extract_text() or "").strip())
        except Exception as e:  # 破損ページで全体を落とさない
            out.append(f"[抽出失敗: {e}]")
    return "\n".join(out)


def extract_text(
    data: bytes,
    filename: str,
    mimetype: str = "",
    max_chars: int = DEFAULT_MAX_CHARS,
) -> dict:
    """Extract plain text. Returns {kind, text, chars, truncated, note}.

    Never raises for an unsupported or corrupt file — it reports the problem in `note`
    so one bad attachment does not abort grading of the rest.
    """
    kind = sniff(filename, mimetype)
    note = ""
    try:
        if kind == "docx":
            text = _docx_text(data)
        elif kind == "pptx":
            text = _pptx_text(data)
        elif kind == "xlsx":
            text = _xlsx_text(data)
        elif kind == "pdf":
            text = _pdf_text(data)
        elif kind == "text":
            text = _decode(data)
        elif kind == "image":
            text = ""
            note = "画像ファイルです。read_submission_images で画像として読み取ってください。"
        else:
            text = ""
            note = f"未対応の形式です（{filename}）。"
    except Exception as e:
        text = ""
        note = f"抽出に失敗しました: {type(e).__name__}: {e}"

    truncated = False
    if max_chars and len(text) > max_chars:
        text = text[:max_chars]
        truncated = True
        note = (note + " " if note else "") + f"{max_chars}字で打ち切りました。"

    if kind in ("docx", "pptx", "xlsx", "pdf") and not text.strip() and not note:
        note = (
            "テキストが取り出せませんでした。画像として貼り付けられている可能性があります。"
            "read_submission_images を試してください。"
        )
    return {"kind": kind, "text": text, "chars": len(text), "truncated": truncated, "note": note}


def _fit_image(raw: bytes, max_px: int) -> tuple[bytes, str] | None:
    """Downscale and re-encode so the image fits model input limits. None if undecodable."""
    from PIL import Image as PILImage

    try:
        im = PILImage.open(io.BytesIO(raw))
        im.load()
    except Exception:
        return None
    if im.mode not in ("RGB", "RGBA", "L"):
        im = im.convert("RGBA" if "A" in im.mode else "RGB")
    if max(im.size) > max_px:
        im.thumbnail((max_px, max_px), PILImage.LANCZOS)

    # スクリーンショットや図は文字が読めることが重要なので、まず可逆の PNG を試す。
    buf = io.BytesIO()
    im.save(buf, format="PNG", optimize=True)
    if buf.tell() <= _IMAGE_TARGET_BYTES:
        return buf.getvalue(), "png"
    buf = io.BytesIO()
    im.convert("RGB").save(buf, format="JPEG", quality=88, optimize=True)
    return buf.getvalue(), "jpeg"


def _zip_media(data: bytes, prefixes: tuple[str, ...]) -> list[tuple[str, bytes]]:
    out: list[tuple[str, bytes]] = []
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        for name in sorted(z.namelist()):
            if name.lower().startswith(prefixes) and not name.endswith("/"):
                out.append((os.path.basename(name), z.read(name)))
    return out


def _pptx_media(data: bytes) -> list[tuple[str, bytes]]:
    """Slide-numbered images, so a grader can tell WHICH slide a figure is on."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    out: list[tuple[str, bytes]] = []
    seen: set[bytes] = set()
    for n, slide in enumerate(prs.slides, 1):
        idx = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or getattr(shape, "image", None):
                try:
                    blob = shape.image.blob
                except Exception:
                    continue
                if blob in seen:
                    continue
                seen.add(blob)
                idx += 1
                out.append((f"slide{n}_img{idx}", blob))
    if not out:  # グループ化された図など python-pptx が拾えない場合の保険
        out = _zip_media(data, ("ppt/media/",))
    return out


def _pdf_media(data: bytes) -> list[tuple[str, bytes]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    out: list[tuple[str, bytes]] = []
    for n, page in enumerate(reader.pages, 1):
        try:
            for i, img in enumerate(page.images, 1):
                out.append((f"page{n}_img{i}", img.data))
        except Exception:
            continue
    return out


def extract_images(
    data: bytes,
    filename: str,
    mimetype: str = "",
    limit: int = 20,
    max_px: int = DEFAULT_MAX_PX,
) -> tuple[list[dict], str]:
    """Extract embedded images. Returns (images, note).

    images: [{name, data: bytes, format: "png"|"jpeg"}] — ready to hand to a model.
    """
    kind = sniff(filename, mimetype)
    raws: list[tuple[str, bytes]] = []
    note = ""
    try:
        if kind == "image":
            raws = [(filename, data)]
        elif kind == "pptx":
            raws = _pptx_media(data)
        elif kind == "docx":
            raws = _zip_media(data, ("word/media/",))
        elif kind == "xlsx":
            raws = _zip_media(data, ("xl/media/",))
        elif kind == "pdf":
            raws = _pdf_media(data)
        else:
            return [], f"この形式からは画像を取り出せません（{filename}）。"
    except Exception as e:
        return [], f"画像の取り出しに失敗しました: {type(e).__name__}: {e}"

    total = len(raws)
    if limit and total > limit:
        raws = raws[:limit]
        note = f"画像は{total}枚あり、先頭{limit}枚のみ返しました。"

    images: list[dict] = []
    skipped = 0
    for name, raw in raws:
        fitted = _fit_image(raw, max_px)
        if fitted is None:
            skipped += 1
            continue
        blob, fmt = fitted
        images.append({"name": name, "data": blob, "format": fmt})
    if skipped:
        note = (note + " " if note else "") + f"{skipped}枚は画像として読めず除外しました。"
    return images, note
